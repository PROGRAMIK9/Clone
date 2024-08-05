from pptx import Presentation
from pptx.util import Inches

# Create a presentation object
prs = Presentation()

# Title Slide
slide_title = prs.slides.add_slide(prs.slide_layouts[0])
title = slide_title.shapes.title
subtitle = slide_title.placeholders[1]

title.text = "Social Media Marketing: Trends and Challenges"
subtitle.text = "An Overview of Recent Developments and Issues"

# Function to add a slide with a title and bullet points
def add_bullet_slide(title, bullet_points):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = title

    tf = body_shape.text_frame
    for point in bullet_points:
        p = tf.add_paragraph()
        p.text = point
        p.level = 0

# Slide 1: Recent Trends in Social Media Marketing
trends = [
    "Video Content Dominance",
    "Influencer Marketing",
    "Social Commerce",
    "Ephemeral Content",
    "Augmented Reality (AR) and Virtual Reality (VR)",
    "User-Generated Content (UGC)",
    "Personalization and Customization",
    "Social Media Stories and Highlights",
    "Purpose-Driven Campaigns",
    "AI and Chatbots"
]
add_bullet_slide("Recent Trends in Social Media Marketing", trends)

# Slide 2: Video Content Dominance
video_trends = [
    "Short-Form Videos: Platforms like TikTok, Instagram Reels, and YouTube Shorts",
    "Live Streaming: Real-time engagement on Facebook, Instagram, and YouTube"
]
add_bullet_slide("Video Content Dominance", video_trends)

# Slide 3: Influencer Marketing
influencer_trends = [
    "Micro-Influencers: Authentic engagement with smaller audiences",
    "Long-Term Partnerships: Sustained relationships over one-off campaigns"
]
add_bullet_slide("Influencer Marketing", influencer_trends)

# Slide 4: Social Commerce
social_commerce_trends = [
    "Shoppable Posts: Direct shopping from Instagram and Facebook posts",
    "In-App Purchases: Seamless shopping experiences within social media platforms"
]
add_bullet_slide("Social Commerce", social_commerce_trends)

# Slide 5: Ephemeral Content
ephemeral_content_trends = [
    "Stories: 24-hour content on Instagram, Snapchat, etc."
]
add_bullet_slide("Ephemeral Content", ephemeral_content_trends)

# Slide 6: AR and VR
ar_vr_trends = [
    "AR Filters: Interactive experiences on Snapchat and Instagram",
    "VR Experiences: Immersive brand experiences and product demonstrations"
]
add_bullet_slide("Augmented Reality (AR) and Virtual Reality (VR)", ar_vr_trends)

# Slide 7: User-Generated Content (UGC)
ugc_trends = [
    "Encouraging UGC: Authentic content created by users",
    "UGC Campaigns: Building community and brand loyalty"
]
add_bullet_slide("User-Generated Content (UGC)", ugc_trends)

# Slide 8: Personalization and Customization
personalization_trends = [
    "Tailored Content: Personalized content based on user data",
    "Dynamic Ads: Changing content based on viewer behavior"
]
add_bullet_slide("Personalization and Customization", personalization_trends)

# Slide 9: Stories and Highlights
stories_highlights_trends = [
    "Stories Highlights: Keeping ephemeral content accessible and organized"
]
add_bullet_slide("Social Media Stories and Highlights", stories_highlights_trends)

# Slide 10: Purpose-Driven Campaigns
purpose_driven_trends = [
    "Social Issues: Aligning with causes that resonate with the audience",
    "Corporate Social Responsibility (CSR): Highlighting CSR efforts"
]
add_bullet_slide("Purpose-Driven Campaigns", purpose_driven_trends)

# Slide 11: AI and Chatbots
ai_chatbots_trends = [
    "Automated Customer Service: 24/7 support on social media",
    "Personalized Recommendations: AI-driven content and product suggestions"
]
add_bullet_slide("AI and Chatbots", ai_chatbots_trends)

# Slide 12: Challenges in Social Media Marketing
challenges = [
    "Algorithm Changes",
    "Data Privacy and Security",
    "Content Saturation",
    "Managing Negative Feedback",
    "Measuring ROI",
    "Keeping Up with Trends",
    "Authenticity and Transparency",
    "Balancing Automation and Human Touch",
    "Platform-Specific Strategies",
    "Crisis Management"
]
add_bullet_slide("Challenges in Social Media Marketing", challenges)

# Slide 13: Algorithm Changes
algorithm_changes_challenges = [
    "Unpredictable Reach: Impact on organic reach and engagement",
    "Adaptation: Constantly adjusting strategies to stay effective"
]
add_bullet_slide("Algorithm Changes", algorithm_changes_challenges)

# Slide 14: Data Privacy and Security
data_privacy_challenges = [
    "Regulations: Compliance with GDPR, CCPA, and other data protection laws",
    "User Trust: Maintaining trust regarding data usage"
]
add_bullet_slide("Data Privacy and Security", data_privacy_challenges)

# Slide 15: Content Saturation
content_saturation_challenges = [
    "High Competition: Difficult to stand out among the vast amount of content",
    "Quality Over Quantity: Creating high-quality content to engage users"
]
add_bullet_slide("Content Saturation", content_saturation_challenges)

# Slide 16: Managing Negative Feedback
negative_feedback_challenges = [
    "Public Criticism: Handling negative comments and reviews",
    "Reputation Management: Quickly resolving issues to maintain brand reputation"
]
add_bullet_slide("Managing Negative Feedback", negative_feedback_challenges)

# Slide 17: Measuring ROI
measuring_roi_challenges = [
    "Attribution: Accurately attributing social media activities to business outcomes",
    "Metrics Overload: Identifying the most relevant metrics to track and analyze"
]
add_bullet_slide("Measuring ROI", measuring_roi_challenges)

# Slide 18: Keeping Up with Trends
keeping_up_with_trends_challenges = [
    "Constant Evolution: Rapid changes in social media trends and features",
    "Platform Updates: Staying updated with new features and algorithm changes"
]
add_bullet_slide("Keeping Up with Trends", keeping_up_with_trends_challenges)

# Slide 19: Authenticity and Transparency
authenticity_challenges = [
    "Trust Issues: Skepticism towards inauthentic content and paid promotions",
    "Genuine Engagement: Creating authentic and transparent content"
]
add_bullet_slide("Authenticity and Transparency", authenticity_challenges)

# Slide 20: Balancing Automation and Human Touch
automation_human_touch_challenges = [
    "Over-Automation: Risk of losing personal touch",
    "Human Interaction: Balancing automated responses with human engagement"
]
add_bullet_slide("Balancing Automation and Human Touch", automation_human_touch_challenges)

# Slide 21: Platform-Specific Strategies
platform_specific_challenges = [
    "Diverse Platforms: Unique audience and best practices for each platform",
    "Tailored Content: Developing specific strategies for each platform"
]
add_bullet_slide("Platform-Specific Strategies", platform_specific_challenges)

# Slide 22: Crisis Management
crisis_management_challenges = [
    "Real-Time Response: Need for quick and effective responses to crises",
    "Preparedness: Having a crisis management plan in place for social media"
]
add_bullet_slide("Crisis Management", crisis_management_challenges)

# Save the presentation
pptx_path = "C:\c programs\social_media.pptx"
prs.save(pptx_path)

pptx_path
